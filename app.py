import streamlit as st
from pptx import Presentation
from moviepy import VideoFileClip
import tempfile
import os

# Create tabs
tab1, tab2 = st.tabs([
    "📄 PPT Notes Extractor",
    "🎥 MP4 to MP3 Extractor"
])

# -------------------------
# TAB 1 - PPT NOTES EXTRACTOR
# -------------------------
with tab1:

    st.title("PowerPoint Notes Extractor")

    uploaded_file = st.file_uploader(
        "Upload a PowerPoint file",
        type=["pptx"],
        key="ppt_upload"
    )

    if uploaded_file:

        with tempfile.NamedTemporaryFile(
            delete=False,
            suffix=".pptx"
        ) as tmp:

            tmp.write(uploaded_file.read())
            ppt_path = tmp.name

        prs = Presentation(ppt_path)

        all_notes = []

        for slide_num, slide in enumerate(prs.slides, start=1):

            if slide.has_notes_slide:

                notes_text = slide.notes_slide.notes_text_frame.text.strip()

                if (
                    notes_text and
                    notes_text.lower() != "click to add notes"
                ):

                    all_notes.append(
                        f"SLIDE {slide_num}\n\n{notes_text}"
                    )

        output_text = "\n\n" + ("-" * 50 + "\n\n").join(all_notes)

        st.success(
            f"Extracted notes from {len(all_notes)} slides."
        )

        st.download_button(
            label="Download TXT File",
            data=output_text,
            file_name="presentation_notes.txt",
            mime="text/plain"
        )

        os.remove(ppt_path)

# -------------------------
# TAB 2 - MP4 TO MP3
# -------------------------
with tab2:

    st.title("MP4 to MP3 Extractor")

    st.write(
        "Upload an MP4 video and download the extracted MP3 audio."
    )

    uploaded_video = st.file_uploader(
        "Upload MP4 Video",
        type=["mp4"],
        key="video_upload"
    )

    if uploaded_video:

        with tempfile.NamedTemporaryFile(
            delete=False,
            suffix=".mp4"
        ) as temp_video:

            temp_video.write(uploaded_video.read())
            video_path = temp_video.name

        mp3_path = video_path.replace(".mp4", ".mp3")

        try:

            video = VideoFileClip(video_path)

            video.audio.write_audiofile(
                mp3_path,
                logger=None
            )

            video.close()

            st.success("Audio extracted successfully!")

            with open(mp3_path, "rb") as f:

                st.download_button(
                    label="Download MP3",
                    data=f,
                    file_name="audio.mp3",
                    mime="audio/mpeg"
                )

        except Exception as e:

            st.error(f"Error: {e}")

        finally:

            if os.path.exists(video_path):
                os.remove(video_path)

            if os.path.exists(mp3_path):
                os.remove(mp3_path)
